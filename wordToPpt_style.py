from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def word_to_ppt(word_path, ppt_template_path, output_ppt_path, font_size=24, font_color=(0, 0, 0)):
    # Word 문서 열기
    doc = Document(word_path)
    prs = Presentation(ppt_template_path)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # RGBColor로 변환
    font_rgb = RGBColor(*font_color)

    text_frame = None  # 본문 텍스트 프레임 초기화

    for para in doc.paragraphs:
        para_text = para.text.strip()  # 공백 제거
        style_name = para.style.name   # 스타일 확인

        if not para_text:
            continue  # 빈 문장은 무시

        if "Heading 1" in style_name:  # 제목 1 → 새로운 슬라이드 생성
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 제목 + 내용 슬라이드 생성
            title = slide.shapes.title

            # 🔹 제목이 없으면 직접 추가
            if title is None:
                title = slide.shapes.add_textbox(Inches(1), Inches(0.5), slide_width - Inches(2), Inches(1))
                title_tf = title.text_frame
            else:
                title_tf = title.text_frame

            title_tf.text = para_text  # 제목 추가

            # 🔹 제목 글꼴 색상 적용
            if title_tf.paragraphs:
                title_tf.paragraphs[0].font.color.rgb = font_rgb

            # 본문을 담을 텍스트 박스 추가 (위치 조정)
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), slide_width - Inches(2), slide_height - Inches(2.5))
            text_frame = textbox.text_frame  # 텍스트 프레임 가져오기
            text_frame.word_wrap = True  # 자동 줄바꿈

        elif "Heading 2" in style_name or "Normal" in style_name:  # 본문 내용
            if slide and text_frame:  # 현재 슬라이드가 있고 본문을 담을 수 있다면
                p = text_frame.add_paragraph()
                p.text = para_text  # 줄바꿈하여 내용 추가
                p.space_after = Pt(10)  # 문단 간격 설정
                p.alignment = PP_ALIGN.LEFT  # 좌측 정렬

                # 글꼴 설정
                run = p.runs[0]
                run.font.size = Pt(font_size)  # 폰트 크기 설정
                run.font.color.rgb = font_rgb  # 🔹 글자 색상 적용

    # 최종 PPT 저장
    prs.save(output_ppt_path)
    print(f"✅ PPT 저장 완료: {output_ppt_path}")

# 🔹 실행 (폰트 크기와 글자 색상을 원하는 값으로 변경 가능)
word_path = "C:/cji_d/test/ppt/test.docx"
ppt_template_path = "C:/cji_d/test/ppt/template/test.pptx"
output_ppt_path = "C:/cji_d/test/ppt/test_output_ppt_path.pptx"

black = (0,0,0)
white = (255,255,255)

# 예제: 글자 크기 20pt, 글자 색상을 검은색으로 설정
word_to_ppt(word_path, ppt_template_path, output_ppt_path, font_size=20, font_color=black)
